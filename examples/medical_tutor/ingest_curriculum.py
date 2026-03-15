#!/usr/bin/env python3
"""Curriculum ingestion pipeline — index medical/nursing content into
both the RAG memory store and the structured knowledge graph.

Supports mixed formats: PDF, DOCX, PPTX, Markdown, plain text.

Usage:
    python examples/medical_tutor/ingest_curriculum.py --help

    # Index a directory of textbooks
    python examples/medical_tutor/ingest_curriculum.py \
        --docs-path ./curriculum/

    # Index with knowledge graph extraction
    python examples/medical_tutor/ingest_curriculum.py \
        --docs-path ./curriculum/ --build-kg

    # Index Open RN textbooks (downloads from NCBI Bookshelf)
    python examples/medical_tutor/ingest_curriculum.py --fetch-open-rn

    # Customize chunking for dense clinical content
    python examples/medical_tutor/ingest_curriculum.py \
        --docs-path ./pharmacology/ \
        --chunk-size 384 --chunk-overlap 96 \
        --subject pharmacology
"""

from __future__ import annotations

import argparse
import json
import os
import sys
import time
import urllib.request
from pathlib import Path
from typing import Any, Dict, List, Optional

# ---------------------------------------------------------------------------
# Open RN textbook catalog (NCBI Bookshelf + WisTech Open)
# ---------------------------------------------------------------------------

OPEN_RN_SOURCES = [
    {
        "title": "Nursing Fundamentals",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK590025/",
        "subject": "fundamentals",
        "description": "Foundational principles of nursing practice",
    },
    {
        "title": "Nursing Pharmacology",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK590024/",
        "subject": "pharmacology",
        "description": "Drug classifications, mechanisms, and nursing implications",
    },
    {
        "title": "Nursing Skills",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK590023/",
        "subject": "skills",
        "description": "Clinical nursing skills and procedures",
    },
    {
        "title": "Nursing Mental Health",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK590029/",
        "subject": "mental_health",
        "description": "Psychiatric and mental health nursing concepts",
    },
    {
        "title": "Health Promotion",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK590028/",
        "subject": "health_promotion",
        "description": "Health assessment, wellness, and disease prevention",
    },
]

# University of Michigan Open Medical Resources
UMICH_SUBJECTS = [
    "anatomy", "physiology", "pathophysiology", "microbiology",
    "biochemistry", "clinical_medicine",
]

# Medical knowledge graph entity types for structured extraction
MEDICAL_ENTITY_TYPES = [
    "disease", "condition", "symptom", "sign", "medication", "drug_class",
    "procedure", "intervention", "nursing_diagnosis", "lab_test",
    "vital_sign", "body_system", "organ", "pathogen", "risk_factor",
    "assessment_finding", "clinical_guideline", "nursing_outcome",
]

MEDICAL_RELATION_TYPES = [
    "causes", "treats", "contraindicated_with", "manifests_as",
    "assessed_by", "indicated_for", "side_effect_of", "interacts_with",
    "part_of", "precedes", "complication_of", "risk_factor_for",
    "nursing_intervention_for", "monitored_by", "classified_as",
]


def _convert_docx(path: Path) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        print(
            f"  Skipping {path.name} — install python-docx: "
            "pip install python-docx",
            file=sys.stderr,
        )
        return ""
    doc = Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _convert_pptx(path: Path) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        print(
            f"  Skipping {path.name} — install python-pptx: "
            "pip install python-pptx",
            file=sys.stderr,
        )
        return ""
    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def read_any_document(path: Path) -> str:
    """Read a document of any supported format and return plain text."""
    suffix = path.suffix.lower()

    if suffix == ".docx":
        return _convert_docx(path)
    elif suffix == ".pptx":
        return _convert_pptx(path)

    # Fall through to OpenJarvis built-in reader for pdf, md, txt, etc.
    try:
        from openjarvis.tools.storage.ingest import read_document
        text, _meta = read_document(path)
        return text
    except ImportError:
        # Fallback: plain text read
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return path.read_text(encoding="latin-1")


def index_documents(
    docs_path: str,
    *,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    subject: str = "general",
    build_kg: bool = False,
    model: str = "qwen3.5:latest",
    engine_key: str = "ollama",
) -> Dict[str, Any]:
    """Index documents into RAG memory and optionally build knowledge graph."""
    from openjarvis import Jarvis
    from openjarvis.tools.storage.chunking import Chunk, ChunkConfig, chunk_text

    j = Jarvis(model=model, engine_key=engine_key)
    stats = {
        "files_processed": 0,
        "chunks_indexed": 0,
        "kg_entities": 0,
        "kg_relations": 0,
        "errors": [],
        "subjects": {},
    }

    root = Path(docs_path)
    if not root.exists():
        print(f"Error: path does not exist: {docs_path}", file=sys.stderr)
        j.close()
        sys.exit(1)

    # Collect all files
    if root.is_file():
        files = [root]
    else:
        files = sorted(root.rglob("*"))
        files = [f for f in files if f.is_file() and not f.name.startswith(".")]

    supported = {
        ".pdf", ".txt", ".md", ".markdown", ".docx", ".pptx",
        ".json", ".csv", ".html", ".htm", ".xml",
    }

    cfg = ChunkConfig(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    total = len([f for f in files if f.suffix.lower() in supported])
    processed = 0

    for filepath in files:
        if filepath.suffix.lower() not in supported:
            continue

        processed += 1
        print(f"  [{processed}/{total}] {filepath.name}...", end=" ", flush=True)

        try:
            text = read_any_document(filepath)
            if not text.strip():
                print("(empty, skipped)")
                continue

            # Chunk the document
            chunks = chunk_text(text, source=str(filepath), config=cfg)

            # Index each chunk into memory with subject metadata
            for chunk in chunks:
                j.memory._get_backend().store(
                    chunk.content,
                    source=chunk.source,
                    metadata={
                        "index": chunk.index,
                        "subject": subject,
                        "file": filepath.name,
                    },
                )

            stats["files_processed"] += 1
            stats["chunks_indexed"] += len(chunks)
            stats["subjects"][subject] = (
                stats["subjects"].get(subject, 0) + len(chunks)
            )
            print(f"({len(chunks)} chunks)")

        except Exception as exc:
            print(f"ERROR: {exc}")
            stats["errors"].append({"file": str(filepath), "error": str(exc)})

    # Build knowledge graph if requested
    if build_kg:
        print("\nBuilding knowledge graph from indexed content...")
        kg_stats = _build_knowledge_graph(j, model, subject)
        stats["kg_entities"] = kg_stats.get("entities", 0)
        stats["kg_relations"] = kg_stats.get("relations", 0)

    j.close()
    return stats


def _build_knowledge_graph(
    j: Any,
    model: str,
    subject: str,
) -> Dict[str, int]:
    """Use the LLM to extract medical entities and relationships from
    indexed content and populate the knowledge graph.
    """
    from openjarvis.tools.storage.knowledge_graph import (
        Entity,
        KnowledgeGraphMemory,
        Relation,
    )

    kg = KnowledgeGraphMemory()
    entity_count = 0
    relation_count = 0

    # Query the memory for representative chunks to extract entities from
    extraction_prompt = f"""Analyze the following medical/nursing text and extract structured entities and relationships.

Return a JSON object with:
- "entities": list of {{"id": "snake_case_id", "type": "<one of: {', '.join(MEDICAL_ENTITY_TYPES)}>", "name": "Display Name", "properties": {{"definition": "...", "system": "body system"}}}}
- "relations": list of {{"source": "entity_id", "target": "entity_id", "type": "<one of: {', '.join(MEDICAL_RELATION_TYPES)}>", "weight": 0.0-1.0}}

Only output valid JSON. Extract key clinical concepts, medications, conditions, and their relationships.

Text:
"""

    # Search for diverse content across the indexed curriculum
    try:
        search_queries = [
            "pathophysiology disease mechanism",
            "medication drug nursing implications",
            "assessment findings vital signs",
            "nursing intervention care plan",
            "laboratory diagnostic test values",
        ]

        for query in search_queries:
            results = j.memory.search(query, top_k=3)
            for result in results:
                content = result.get("content", "")
                if len(content) < 100:
                    continue

                try:
                    response = j.ask(
                        extraction_prompt + content[:2000],
                        temperature=0.1,
                        context=False,
                    )

                    # Parse the JSON response
                    # Try to extract JSON from the response
                    json_str = response
                    if "```json" in response:
                        json_str = response.split("```json")[1].split("```")[0]
                    elif "```" in response:
                        json_str = response.split("```")[1].split("```")[0]

                    data = json.loads(json_str.strip())

                    for ent in data.get("entities", []):
                        entity = Entity(
                            entity_id=ent["id"],
                            entity_type=ent["type"],
                            name=ent["name"],
                            properties={
                                **ent.get("properties", {}),
                                "subject": subject,
                            },
                        )
                        kg.add_entity(entity)
                        entity_count += 1

                    for rel in data.get("relations", []):
                        relation = Relation(
                            source_id=rel["source"],
                            target_id=rel["target"],
                            relation_type=rel["type"],
                            weight=rel.get("weight", 1.0),
                        )
                        kg.add_relation(relation)
                        relation_count += 1

                except (json.JSONDecodeError, KeyError, IndexError):
                    continue

    except Exception as exc:
        print(f"  Knowledge graph extraction error: {exc}", file=sys.stderr)

    kg.close()
    return {"entities": entity_count, "relations": relation_count}


def fetch_open_rn(download_dir: str = "./curriculum/open_rn") -> List[str]:
    """Provide instructions for obtaining Open RN textbooks.

    The Open RN textbooks are CC-BY 4.0 licensed and available from:
    - NCBI Bookshelf: https://www.ncbi.nlm.nih.gov/books/NBK590025/
    - WisTech Open: https://www.wistechopen.org/open-rn-details

    Returns a list of source descriptions.
    """
    dest = Path(download_dir)
    dest.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("  Open RN Nursing Textbooks (CC-BY 4.0)")
    print("=" * 60)
    print()
    print("The Open RN project provides free, peer-reviewed nursing")
    print("textbooks. Download them and place in the curriculum directory.")
    print()
    print("Available textbooks:")
    print()

    sources = []
    for src in OPEN_RN_SOURCES:
        print(f"  - {src['title']}")
        print(f"    Subject: {src['subject']}")
        print(f"    URL: {src['url']}")
        print(f"    Description: {src['description']}")
        print()
        sources.append(src["title"])

    print("Additional resources:")
    print()
    print("  - University of Michigan Open Medical Resources")
    print("    https://open.umich.edu/find/open-educational-resources/medical-resources")
    print()
    print("  - WisTech Open (full catalog)")
    print("    https://www.wistechopen.org/open-rn-details")
    print()

    # Create a README in the download directory
    readme = dest / "README.md"
    readme.write_text(
        "# Medical Curriculum Sources\n\n"
        "Place your curriculum files in this directory, organized by subject:\n\n"
        "```\n"
        "curriculum/\n"
        "├── open_rn/\n"
        "│   ├── fundamentals/\n"
        "│   ├── pharmacology/\n"
        "│   ├── skills/\n"
        "│   ├── mental_health/\n"
        "│   └── health_promotion/\n"
        "├── anatomy/\n"
        "├── pathophysiology/\n"
        "├── clinical_guidelines/\n"
        "└── board_prep/\n"
        "```\n\n"
        "## Open Educational Resources\n\n"
        "### Open RN (CC-BY 4.0)\n"
        "- NCBI Bookshelf: https://www.ncbi.nlm.nih.gov/books/NBK590025/\n"
        "- WisTech Open: https://www.wistechopen.org/open-rn-details\n\n"
        "### University of Michigan Medical OER\n"
        "- https://open.umich.edu/find/open-educational-resources/medical-resources\n\n"
        "## Supported Formats\n\n"
        "- PDF (.pdf) — textbooks, clinical guidelines\n"
        "- Word (.docx) — lecture notes, study guides\n"
        "- PowerPoint (.pptx) — lecture slides\n"
        "- Markdown (.md) — notes, summaries\n"
        "- Text (.txt) — plain text content\n"
        "- HTML (.html) — web-scraped content\n"
    )
    print(f"Created {readme}")
    print()
    print("After downloading, run:")
    print(f"  python examples/medical_tutor/ingest_curriculum.py \\")
    print(f"      --docs-path {download_dir} --build-kg")

    return sources


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Index medical/nursing curriculum into OpenJarvis",
    )
    parser.add_argument(
        "--docs-path",
        type=str,
        help="Path to curriculum directory or file to index.",
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=512,
        help="Chunk size in tokens (default: 512).",
    )
    parser.add_argument(
        "--chunk-overlap",
        type=int,
        default=64,
        help="Overlap between chunks (default: 64).",
    )
    parser.add_argument(
        "--subject",
        type=str,
        default="general",
        help=(
            "Subject tag for organization "
            "(e.g., pharmacology, fundamentals, pathophysiology)."
        ),
    )
    parser.add_argument(
        "--build-kg",
        action="store_true",
        help="Extract entities/relations into the knowledge graph.",
    )
    parser.add_argument(
        "--model",
        type=str,
        default="qwen3.5:latest",
        help="Model for KG extraction (default: qwen3.5:latest).",
    )
    parser.add_argument(
        "--engine",
        type=str,
        default="ollama",
        help="Inference engine (default: ollama).",
    )
    parser.add_argument(
        "--fetch-open-rn",
        action="store_true",
        help="Show instructions for obtaining Open RN textbooks.",
    )
    args = parser.parse_args()

    if args.fetch_open_rn:
        fetch_open_rn()
        return

    if not args.docs_path:
        parser.error("--docs-path is required (or use --fetch-open-rn)")

    try:
        from openjarvis import Jarvis  # noqa: F401
    except ImportError:
        print(
            "Error: openjarvis is not installed. "
            "Install with: uv sync --extra dev",
            file=sys.stderr,
        )
        sys.exit(1)

    print("=" * 60)
    print("  Medical Curriculum Ingestion Pipeline")
    print("=" * 60)
    print()
    print(f"Source: {args.docs_path}")
    print(f"Subject: {args.subject}")
    print(f"Chunk size: {args.chunk_size} (overlap: {args.chunk_overlap})")
    print(f"Knowledge graph: {'yes' if args.build_kg else 'no'}")
    print(f"Model: {args.model} ({args.engine})")
    print("-" * 60)
    print()

    t0 = time.time()
    stats = index_documents(
        args.docs_path,
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
        subject=args.subject,
        build_kg=args.build_kg,
        model=args.model,
        engine_key=args.engine,
    )
    elapsed = time.time() - t0

    print()
    print("=" * 60)
    print("  Ingestion Complete")
    print("=" * 60)
    print(f"  Files processed:  {stats['files_processed']}")
    print(f"  Chunks indexed:   {stats['chunks_indexed']}")
    if args.build_kg:
        print(f"  KG entities:      {stats['kg_entities']}")
        print(f"  KG relations:     {stats['kg_relations']}")
    print(f"  Errors:           {len(stats['errors'])}")
    print(f"  Time:             {elapsed:.1f}s")
    print()

    if stats["subjects"]:
        print("  Chunks by subject:")
        for subj, count in stats["subjects"].items():
            print(f"    {subj}: {count}")
        print()

    if stats["errors"]:
        print("  Errors:")
        for err in stats["errors"]:
            print(f"    {err['file']}: {err['error']}")
        print()

    print("Next steps:")
    print("  python examples/medical_tutor/tutor.py")


if __name__ == "__main__":
    main()
